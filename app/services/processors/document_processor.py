"""Document processor for PDF, DOCX, DOC, PPTX, TXT, MD files"""
import tempfile
import os
import logging
from typing import List, Dict, Any
from app.services.processors.base import BaseProcessor, ProcessedChunk, ProcessedContent
from app.services.chunking import chunk_text_custom

logger = logging.getLogger(__name__)


class DocumentProcessor(BaseProcessor):
    """Processor for document files (PDF, DOCX, DOC, PPTX, TXT, MD)"""
    
    def __init__(self):
        self.supported_mime_types = [
            "application/pdf",
            "application/msword",  # .doc
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
            "text/plain",  # .txt
            "text/markdown",  # .md
        ]
        self.supported_extensions = [".pdf", ".doc", ".docx", ".pptx", ".txt", ".md"]
    
    def can_process(self, mime_type: str) -> bool:
        """Check if this processor can handle the MIME type"""
        return mime_type in self.supported_mime_types
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file extensions"""
        return self.supported_extensions
    
    async def process(self, file_content: bytes, filename: str) -> ProcessedContent:
        """
        Process document file and extract text chunks
        
        For PDF/DOCX/DOC: Uses unstructured.partition() and chunk_elements()
        For PPTX: Uses python-pptx to extract text from slides
        For TXT/MD: Direct decode + custom chunking
        """
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext in [".pdf", ".docx", ".doc"]:
            return await self._process_with_unstructured(file_content, filename, file_ext)
        elif file_ext == ".pptx":
            return await self._process_pptx(file_content, filename)
        elif file_ext in [".txt", ".md"]:
            return await self._process_text_file(file_content, filename)
        else:
            raise ValueError(f"Unsupported document format: {file_ext}")
    
    async def _process_with_unstructured(
        self, file_content: bytes, filename: str, file_ext: str
    ) -> ProcessedContent:
        """Process PDF/DOCX/DOC using unstructured library"""
        try:
            from unstructured.partition.auto import partition
            from unstructured.chunking.basic import chunk_elements
            from unstructured.cleaners.core import clean
            
            # Write to temp file (unstructured requires file path)
            with tempfile.NamedTemporaryFile(
                delete=False, suffix=file_ext, mode="wb"
            ) as temp_file:
                temp_file.write(file_content)
                temp_path = temp_file.name
            
            try:
                # Extract elements using unstructured
                elements = partition(filename=temp_path)
                
                # Clean text from elements while preserving metadata
                cleaned_elements = []
                for element in elements:
                    if hasattr(element, "text") and element.text:
                        cleaned_text = clean(element.text)
                        if cleaned_text.strip():
                            # Update element's text with cleaned version (preserves metadata)
                            element.text = cleaned_text
                            cleaned_elements.append(element)
                
                # Chunk elements using unstructured (max_characters=2500, overlap=100, overlap_all=True)
                chunked_elements = chunk_elements(
                    elements=cleaned_elements,
                    max_characters=2500,
                    overlap=100,
                    overlap_all=True  # Apply overlap between all chunks, not just oversized elements
                )
                
                # Extract text from chunked elements
                chunks = []
                full_text_parts = []
                
                for idx, chunked_element in enumerate(chunked_elements):
                    chunk_text = chunked_element.text if hasattr(chunked_element, "text") else str(chunked_element)
                    full_text_parts.append(chunk_text)
                    
                    chunks.append(ProcessedChunk(
                        content=chunk_text,
                        chunk_type="text",
                        chunk_index=idx,
                        metadata={}
                    ))
                
                full_text = "\n\n".join(full_text_parts)
                
                return ProcessedContent(
                    text=full_text,
                    chunks=chunks,
                    metadata={"filename": filename, "method": "unstructured"},
                    structure={
                        "chunk_count": len(chunks),
                        "method": "unstructured.chunk_elements",
                        "max_characters": 2500,
                        "overlap": 100
                    }
                )
            
            finally:
                # Clean up temp file
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
        
        except ImportError:
            logger.error("unstructured library not installed. Install with: pip install 'unstructured[pdf,docx]'")
            raise
        except Exception as e:
            logger.error(f"Error processing document with unstructured: {e}")
            raise
    
    async def _process_text_file(self, file_content: bytes, filename: str) -> ProcessedContent:
        """Process TXT/MD files with direct decode + custom chunking"""
        try:
            # Decode bytes to UTF-8 string
            text = file_content.decode("utf-8")
            
            # Use custom chunking (2500 chars, 100 overlap)
            chunk_texts = chunk_text_custom(text, max_chars=2500, overlap=100)
            
            # Create ProcessedChunk objects
            chunks = []
            for idx, chunk_text in enumerate(chunk_texts):
                chunks.append(ProcessedChunk(
                    content=chunk_text,
                    chunk_type="text",
                    chunk_index=idx,
                    metadata={}
                ))
            
            return ProcessedContent(
                text=text,
                chunks=chunks,
                metadata={"filename": filename, "method": "direct_decode"},
                structure={
                    "chunk_count": len(chunks),
                    "method": "custom_chunking",
                    "max_characters": 2500,
                    "overlap": 100
                }
            )
        
        except UnicodeDecodeError as e:
            logger.error(f"Error decoding text file {filename}: {e}")
            raise ValueError(f"Failed to decode file as UTF-8: {e}")
    
    async def _process_pptx(self, file_content: bytes, filename: str) -> ProcessedContent:
        """Process PPTX files using python-pptx to extract text from slides"""
        try:
            from pptx import Presentation
            import io
            
            # Load presentation from bytes
            presentation = Presentation(io.BytesIO(file_content))
            
            # Extract text from all slides
            full_text_parts = []
            slide_texts = []
            
            for slide_idx, slide in enumerate(presentation.slides):
                slide_text_parts = []
                
                # Extract text from all shapes on the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_text_parts.append(text)
                            full_text_parts.append(text)
                
                # Combine all text from this slide
                if slide_text_parts:
                    slide_text = "\n".join(slide_text_parts)
                    slide_texts.append(slide_text)
            
            # If no text found, return empty content
            if not full_text_parts:
                logger.warning(f"No text extracted from PPTX file {filename}")
                return ProcessedContent(
                    text="",
                    chunks=[],
                    metadata={"filename": filename, "method": "pptx"},
                    structure={
                        "chunk_count": 0,
                        "method": "python-pptx",
                        "slide_count": len(presentation.slides)
                    }
                )
            
            # Combine all text
            full_text = "\n\n".join(full_text_parts)
            
            # Use custom chunking (2500 chars, 100 overlap)
            chunk_texts = chunk_text_custom(full_text, max_chars=2500, overlap=100)
            
            # Create ProcessedChunk objects
            chunks = []
            for idx, chunk_text in enumerate(chunk_texts):
                chunks.append(ProcessedChunk(
                    content=chunk_text,
                    chunk_type="text",
                    chunk_index=idx,
                    metadata={}
                ))
            
            return ProcessedContent(
                text=full_text,
                chunks=chunks,
                metadata={"filename": filename, "method": "pptx"},
                structure={
                    "chunk_count": len(chunks),
                    "method": "python-pptx",
                    "slide_count": len(presentation.slides),
                    "max_characters": 2500,
                    "overlap": 100
                }
            )
        
        except ImportError:
            logger.error("python-pptx library not installed. Install with: pip install python-pptx")
            raise
        except Exception as e:
            logger.error(f"Error processing PPTX file {filename}: {e}")
            raise
